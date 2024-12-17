from pptx import Presentation



import os

import platform



presentation = Presentation()




def add_slide(title, content):
    slide_layout = presentation.slide_layouts[1]

    slide = presentation.slides.add_slide(slide_layout)

    title_box = slide.shapes.title

    content_box = slide.shapes.placeholders[1]

    title_box.text = title

    content_box.text = content


# Slide 1: Title Slide

title_slide_layout = presentation.slide_layouts[0]  # Title Slide

title_slide = presentation.slides.add_slide(title_slide_layout)

title = title_slide.shapes.title

subtitle = title_slide.placeholders[1]

title.text = "Effective Work Tracking: Microsoft Planner and Jira"

subtitle.text = "Exploring Two Powerful Tools for Project Management"

# Slide 2: Introduction

add_slide("Introduction",

          "The importance of tracking work effectively in managing projects.\n"

          "Overview of two tools: Microsoft Planner and Jira, focusing on their strengths and use cases.")

# Slide 3: What is Microsoft Planner?

add_slide("What is Microsoft Planner?",

          "Microsoft Planner is a task management tool integrated within Microsoft 365.\n"

          "Ideal for teams looking for a user-friendly interface for collaboration.")

# Slide 4: Key Features of Microsoft Planner

add_slide("Key Features of Microsoft Planner",

          " Task Management: Assign tasks, set due dates, and track progress.\n"

          " Visual Kanban Boards: Organize work visually using customizable boards.\n"

          " Integration: Works seamlessly with other Microsoft tools like Teams and Outlook.\n"

          " Comments and Attachments: Collaborate within tasks by adding comments and files.")

# Slide 5: Benefits of Using Microsoft Planner

add_slide("Benefits of Using Microsoft Planner",

          " Ease of Use: Intuitive design makes it accessible for all team members.\n"

          " Collaboration-Friendly: Real-time updates and notifications facilitate communication.\n"

          " Cost-Effective: Available as part of Microsoft 365 subscriptions.")

# Slide 6: What is Jira?

add_slide("What is Jira?",

          "Jira is a project management tool developed by Atlassian, primarily used by software development teams.\n"

          "Supports Agile project management methodologies like Scrum and Kanban.")

# Slide 7: Key Features of Jira

add_slide("Key Features of Jira",

          " Agile Boards: Visualize workflows with customizable Scrum and Kanban boards.\n"

          " Advanced Reporting: Generate detailed reports on project progress and team performance.\n"

          " Custom Workflows: Tailor project workflows to meet specific needs.\n"

          " Integration with Development Tools: Connects with tools for bug tracking and code versioning.")

# Slide 8: Benefits of Using Jira

add_slide("Benefits of Using Jira",

          " Designed for Development Teams: Tailored features support software project management.\n"

          " Scalability: Suitable for both small teams and large organizations with complex workflows.\n"

          " Strong Community Support: Extensive resources, forums, and documentation are available.")

# Slide 9: When to Use Microsoft Planner

add_slide("When to Use Microsoft Planner",

          " Best for non-technical teams or organizations using Microsoft 365.\n"

          " Suitable for straightforward project management and collaboration.\n"

          " Ideal for managing tasks, deadlines, and team assignments in a visual format.")

# Slide 10: When to Use Jira

add_slide("When to Use Jira",

          " Best for software development and technical teams.\n"

          " Ideal for projects that require agile methodologies and detailed reporting.\n"

          " Can accommodate complex workflows and integrate with development tools.")

# Slide 11: Summary

add_slide("Summary",

          " Microsoft Planner provides a simple, cost-effective solution for general project management.\n"

          " Jira offers robust features for teams working on software development, making it a strong choice for agile methodologies.\n"

          " Both tools can enhance productivity and organization depending on team needs and project requirements.")

# Slide 12: Questions?

add_slide("Questions?", "Open the floor for questions and encourage discussion on experiences with either tool.")

# Slide 13: References

add_slide("References",

          " Microsoft website and Atlassian resources for further reading.")

# Save the presentation

pptx_file = "Work_Tracking_Presentation.pptx"

presentation.save(pptx_file)

# Open the presentation

platform.system()
os.startfile(pptx_file)



print("Presentation created and opened successfully!")